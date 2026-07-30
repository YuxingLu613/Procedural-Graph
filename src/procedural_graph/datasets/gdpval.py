"""GDPval evaluation environment adapter subclassing BaseEnvironment."""

import os
import re
import shutil
import subprocess
import tempfile
import json
import zipfile
import xml.etree.ElementTree as ET
import getpass
from typing import Any, Dict, List, Optional, Tuple
import pandas as pd
from .. import env_base

# ====================================================================
# LLM Judge Prompts
# ====================================================================

JUDGE_PROMPT_TEMPLATE = """
You are a professional evaluator. Your task is to verify whether the agent's output artifacts satisfy the criteria provided in the rubrics.

Here are the agent's output artifacts (and possibly reference files):
{artifacts_content}

Here is the rubrics JSON containing the criteria to verify:
{rubrics_content}

For each criterion, evaluate whether it is satisfied (true) or not (false) based on the output artifacts.
Provide a clear justification and list the sources (filenames) you referenced.

After all verifications complete, output the final ratings as a JSON array containing all criteria with their verification results. Each item in the array MUST be a JSON object with a SINGLE KEY named "criterion N" (where N is the criterion number, starting from 1).
The value of that key must be a JSON object with the following fields:
- "description": string with the criterion description
- "satisfied": boolean (true or false)
- "justification": string with your reasoning
- "sources": list of filenames you referenced to verify the criterion

**CRITICAL**: You MUST wrap your final JSON output in <FINAL_RATINGS></FINAL_RATINGS> XML tags. This is REQUIRED for parsing. Your response MUST end with these tags.

CORRECT format:
<FINAL_RATINGS>
[
  {{"criterion 1": {{"description": "...", "satisfied": true, "justification": "...", "sources": [...]}}}},
  {{"criterion 2": {{"description": "...", "satisfied": false, "justification": "...", "sources": [...]}}}},
  ...
]
</FINAL_RATINGS>

Do NOT output any text after the closing </FINAL_RATINGS> tag.
"""


def get_splits(df: pd.DataFrame) -> Dict[str, pd.DataFrame]:
  """Deterministically splits the GDPval dataset into Train (88)/Val (88)/Test (44).

  Splits each of the 44 occupations in a 2:2:1 ratio (2 train, 2 val, 1 test).
  """
  train_rows = []
  val_rows = []
  test_rows = []

  # Sort occupations deterministically
  for _, group in sorted(df.groupby("occupation"), key=lambda x: x[0]):
    group = group.sort_values("task_id")
    train_rows.append(group.iloc[:2])
    val_rows.append(group.iloc[2:4])
    test_rows.append(group.iloc[4:5])

  train_df = pd.concat(train_rows).reset_index(drop=True)
  val_df = pd.concat(val_rows).reset_index(drop=True)
  test_df = pd.concat(test_rows).reset_index(drop=True)

  return {
      "train": train_df,
      "val": val_df,
      "test": test_df,
  }





class GDPvalEnv(env_base.BaseEnvironment):
  """GDPval evaluation environment adapter."""

  _df_cache = {}
  _splits_cache = {}
  _judge_client = None

  def __init__(
      self,
      sample_index: int = 0,
      dataset_path: Optional[str] = None,
      llm_client: Optional[Any] = None,
      split: str = "test",
  ):
    super().__init__()
    self.dataset_path = dataset_path or os.path.abspath(os.path.join(os.path.dirname(__file__), "../../../data/datasets/General/GDPval/data/train-00000-of-00001.parquet"))
    self.llm_client = llm_client
    self.split = split

    # Load dataset and cache it
    if self.dataset_path not in GDPvalEnv._df_cache:
      if not os.path.exists(self.dataset_path):
        raise FileNotFoundError(
            f"GDPval dataset file not found at: {self.dataset_path}"
        )
      df = pd.read_parquet(self.dataset_path)
      GDPvalEnv._df_cache[self.dataset_path] = df
      GDPvalEnv._splits_cache[self.dataset_path] = get_splits(df)

    self.splits = GDPvalEnv._splits_cache[self.dataset_path]
    if self.split not in self.splits:
      raise ValueError(f"Invalid split: {self.split}")

    self.df = self.splits[self.split]

    if sample_index < 0 or sample_index >= len(self.df):
      raise IndexError(
          f"sample_index {sample_index} out of range for split '{self.split}'"
          f" (total rows: {len(self.df)})"
      )

    row = self.df.iloc[sample_index]
    self.task_id = row["task_id"]
    self.sector = row["sector"]
    self.occupation = row["occupation"]
    self.prompt = row["prompt"]
    self.target_question = self.prompt
    self.reference_files = (
        list(row["reference_files"])
        if row["reference_files"] is not None
        else []
    )
    self.deliverable_files = (
        list(row["deliverable_files"])
        if row["deliverable_files"] is not None
        else []
    )
    self.rubric_json = row["rubric_json"]

    self.terminated = False
    self.score = 0.0
    self.trajectory_history: List[str] = []
    self.sandbox_dir = ""

  def get_system_prompt(self) -> str:
    return (
        f"You are a professional {self.occupation} working in the {self.sector} sector.\n"
        f"Your task is: {self.prompt}\n\n"
        "You have access to a local environment and specialized tools. "
        "Any files you need to modify or create must be done in the current directory.\n\n"
        "Available Tools:\n"
        " - list_dir(path: str = '.'): List files and directories in the workspace.\n"
        " - view_file(path: str): View the contents of a text file.\n"
        " - read_excel(path: str, sheet_name: Optional[str] = None): Read sheet names or table data from an Excel file.\n"
        " - read_word(path: str): Read text and tables from a Word (.docx) document.\n"
        " - read_pdf(path: str): Extract text content from a PDF (.pdf) file.\n"
        " - search_files(query: str, path: str = '.'): Search for keyword in text files.\n"
        " - run_python(script_content: str): Directly execute a multi-line Python script without bash string escaping. PREFER THIS over code_exec for python scripts.\n"
        " - write_file(path: str, content: str): Write text content to a file. Overwrites if exists.\n"
        " - code_exec(command: str): Execute a bash command in the workspace. Timeout is 60s.\n"
        " - submit(): Submit the final deliverables and end the task. Call this ONLY when you have verified all deliverables are correctly generated.\n\n"
        "Once you have generated all the required deliverables, call the `submit` tool."
    )

  def is_terminated(self) -> bool:
    return self.terminated

  def get_score(self) -> float:
    return self.score

  def reset(self) -> str:
    self.terminated = False
    self.score = 0.0
    self.trajectory_history = []

    # Initialize sandbox in /tmp to avoid CitC permission issues
    # Add PID and short UUID to prevent fatal sandbox collisions during parallel experiments
    username = getpass.getuser()
    import uuid
    self.sandbox_dir = os.path.join(
        tempfile.gettempdir(), f"gdpval_sandbox_{username}_{self.task_id}_{self.split}_{os.getpid()}_{uuid.uuid4().hex[:6]}"
    )


    if os.path.exists(self.sandbox_dir):
      shutil.rmtree(self.sandbox_dir)
    os.makedirs(self.sandbox_dir)

    # Create symlink from python -> /usr/bin/python3 so invoking 'python' works seamlessly
    try:
      if os.path.exists("/usr/bin/python3"):
        os.symlink("/usr/bin/python3", os.path.join(self.sandbox_dir, "python"))
    except Exception:
      pass

    # Copy reference files
    # We resolve their paths relative to the dataset directory
    dataset_dir = os.path.dirname(os.path.dirname(self.dataset_path))
    if self.reference_files is not None:
      for ref_file in self.reference_files:
        src_path = os.path.join(dataset_dir, ref_file)
        dst_path = os.path.join(self.sandbox_dir, os.path.basename(ref_file))
        if os.path.exists(src_path):
          shutil.copy(src_path, dst_path)
        else:
          print(f"Warning: Reference file not found: {src_path}")

    return f"Environment reset. Task: {self.prompt}"

  def get_state(self) -> Any:
    return (self.terminated, self.score, list(self.trajectory_history))

  def load_state(self, state: Any) -> None:
    self.terminated, self.score, traj = state
    self.trajectory_history = list(traj)

  # =========================================================================
  # Registered Tools
  # =========================================================================

  def _get_safe_path(self, path: str) -> str:
    abs_path = os.path.abspath(os.path.join(self.sandbox_dir, path))
    if not abs_path.startswith(self.sandbox_dir):
      raise ValueError("Access denied: Path escapes sandbox.")
    return abs_path

  def _is_binary(self, path: str) -> bool:
    ext = os.path.splitext(path)[1].lower()
    if ext in (
        ".xlsx",
        ".xls",
        ".pdf",
        ".docx",
        ".doc",
        ".zip",
        ".tar",
        ".gz",
        ".png",
        ".jpg",
        ".jpeg",
    ):
      return True
    try:
      with open(path, "rb") as f:
        chunk = f.read(1024)
        return b"\x00" in chunk
    except Exception:
      return False

  @env_base.register_tool("list_dir", "List files in the workspace.")
  def list_dir(self, path: str = ".") -> str:
    try:
      safe_path = self._get_safe_path(path)
      files = os.listdir(safe_path)
      return json.dumps(files)
    except Exception as e:
      return f"Error: {e}"

  @env_base.register_tool("view_file", "View contents of a file.")
  def view_file(self, path: str) -> str:
    try:
      safe_path = self._get_safe_path(path)
      if self._is_binary(safe_path):
        return (
            f"Warning: '{path}' is a binary file. You cannot view it directly"
            " as text. Use read_excel, read_word, read_pdf, or run_python to inspect it."
        )
      with open(safe_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()
    except Exception as e:
      return f"Error: {e}"

  @env_base.register_tool("write_file", "Write content to a file.")
  def write_file(self, path: str, content: str) -> str:
    try:
      safe_path = self._get_safe_path(path)
      os.makedirs(os.path.dirname(safe_path), exist_ok=True)
      with open(safe_path, "w", encoding="utf-8") as f:
        f.write(content)
      return f"Successfully wrote to {path}"
    except Exception as e:
      return f"Error: {e}"

  @env_base.register_tool(
      "code_exec",
      "Execute a bash command in the workspace.",
      is_state_changing=True,
  )
  def code_exec(self, command: str) -> str:
    try:
      env = os.environ.copy()
      env["PATH"] = f"{self.sandbox_dir}:{env.get('PATH', '')}"
      result = subprocess.run(
          command,
          shell=True,
          cwd=self.sandbox_dir,
          env=env,
          stdout=subprocess.PIPE,
          stderr=subprocess.PIPE,
          text=True,
          timeout=60,
      )

      output = f"Exit code: {result.returncode}\n"
      if result.stdout:
        output += f"Stdout:\n{result.stdout}\n"
      if result.stderr:
        output += f"Stderr:\n{result.stderr}\n"
      return output
    except subprocess.TimeoutExpired:
      return "Error: Command timed out after 60 seconds."
    except Exception as e:
      return f"Error: {e}"

  @env_base.register_tool(
      "run_python",
      "Directly execute a multi-line Python script without bash escaping.",
      is_state_changing=True,
  )
  def run_python(self, script_content: str) -> str:
    try:
      script_path = os.path.join(self.sandbox_dir, "_temp_script.py")
      with open(script_path, "w", encoding="utf-8") as f:
        f.write(script_content)
      result = subprocess.run(
          ["/usr/bin/python3", "_temp_script.py"],
          cwd=self.sandbox_dir,
          stdout=subprocess.PIPE,
          stderr=subprocess.PIPE,
          text=True,
          timeout=60,
      )
      output = f"Exit code: {result.returncode}\n"
      if result.stdout:
        output += f"Stdout:\n{result.stdout}\n"
      if result.stderr:
        output += f"Stderr:\n{result.stderr}\n"
      return output
    except subprocess.TimeoutExpired:
      return "Error: Python script timed out after 60 seconds."
    except Exception as e:
      return f"Error executing Python script: {e}"

  @env_base.register_tool(
      "read_excel",
      "Read sheet names or table data from an Excel file (.xlsx/.xls).",
  )
  def read_excel(self, path: str, sheet_name: Optional[str] = None) -> str:
    try:
      safe_path = self._get_safe_path(path)
      xls = pd.ExcelFile(safe_path)
      if sheet_name is None:
        return f"Sheets in {path}: {list(xls.sheet_names)}. Specify `sheet_name` to view content."
      if sheet_name not in xls.sheet_names:
        return f"Error: Sheet '{sheet_name}' not found. Available sheets: {list(xls.sheet_names)}"
      df = pd.read_excel(safe_path, sheet_name=sheet_name)
      return f"### Sheet: {sheet_name}\nShape: {df.shape}\n\n" + df.head(50).to_markdown(index=False)
    except Exception as e:
      return f"Error reading Excel file {path}: {e}"

  @env_base.register_tool(
      "read_word",
      "Read paragraphs and tables from a Word document (.docx).",
  )
  def read_word(self, path: str) -> str:
    try:
      safe_path = self._get_safe_path(path)
      with zipfile.ZipFile(safe_path) as docx:
        xml_content = docx.read("word/document.xml")
        root = ET.fromstring(xml_content)
        namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        text_elems = root.findall(".//w:t", namespaces)
        return "".join([elem.text for elem in text_elems if elem.text])
    except Exception as e:
      return f"Error reading Word document {path}: {e}"

  @env_base.register_tool(
      "read_pdf",
      "Extract text content from a PDF file (.pdf).",
  )
  def read_pdf(self, path: str) -> str:
    try:
      safe_path = self._get_safe_path(path)
      res = subprocess.run(
          ["pdftotext", safe_path, "-"],
          capture_output=True,
          text=True,
          timeout=15,
      )
      if res.returncode == 0 and res.stdout.strip():
        return res.stdout
      return "Warning: pdftotext extracted empty content or failed. Try processing via Python script."
    except Exception as e:
      return f"Error reading PDF {path}: {e}"

  @env_base.register_tool(
      "search_files",
      "Search for a keyword or regex pattern in text files within the workspace.",
  )
  def search_files(self, query: str, path: str = ".") -> str:
    try:
      safe_path = self._get_safe_path(path)
      res = subprocess.run(
          ["grep", "-rnI", "--max-count=20", query, safe_path],
          capture_output=True,
          text=True,
          timeout=15,
      )
      if res.stdout:
        return res.stdout
      return "No matches found."
    except Exception as e:
      return f"Error searching files: {e}"

  @env_base.register_tool(
      "submit",
      "Submit the final deliverables and end the task.",
      is_state_changing=True,
  )
  def submit(self) -> str:
    self.terminated = True
    self.score = self.evaluate()
    return f"Task submitted. Evaluation score: {self.score}"


  # =========================================================================
  # Evaluation (LLM Judge)
  # =========================================================================

  def _extract_text_from_file(self, filepath: str) -> str:
    """Robustly extracts text from various file formats with token length capping."""
    try:
      file_size = os.path.getsize(filepath)
      if file_size > 10 * 1024 * 1024:  # 10MB limit
        return (
            f"[File '{os.path.basename(filepath)}' skipped: exceeds 10MB limit]"
        )
    except Exception:
      pass

    ext = os.path.splitext(filepath)[1].lower()
    max_chars = 150_000  # ~35,000 tokens cap per file to prevent 1M token overflow

    if ext in (".xlsx", ".xls"):
      try:
        xls = pd.ExcelFile(filepath)
        content = ""
        for sheet_name in xls.sheet_names:
          df = pd.read_excel(filepath, sheet_name=sheet_name)
          content += f"### Sheet: {sheet_name}\n"
          if not df.empty:
            # Cap rows if dataframe is massive
            if len(df) > 500:
              content += df.head(500).to_string(index=False) + "\n...[ROWS TRUNCATED]\n\n"
            else:
              content += df.to_string(index=False) + "\n\n"
          else:
            content += "[Empty Sheet]\n\n"
          if len(content) > max_chars:
            return content[:max_chars] + "\n...[TRUNCATED: EXCEL CONTENT EXCEEDS 150,000 CHARACTERS]..."
        return content
      except Exception as e:
        return f"[Error reading Excel file: {e}]"

    elif ext in (".pptx", ".ppt"):
      try:
        from pptx import Presentation
        prs = Presentation(filepath)
        content = ""
        for i, slide in enumerate(prs.slides):
          content += f"### Slide {i+1}\n"
          for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
              content += shape.text + "\n"
          content += "\n"
        if len(content) > max_chars:
          return content[:max_chars] + "\n...[TRUNCATED: PPTX EXCEEDS 150,000 CHARACTERS]..."
        return content
      except Exception as e:
        return f"[Error reading PPTX file: {e}]"

    binary_extensions = {
        ".zip", ".tar", ".gz", ".rar", ".7z",
        ".wav", ".mp3", ".ogg", ".flac", ".m4a",
        ".mp4", ".avi", ".mkv", ".mov",
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".ico",
        ".doc",
    }
    if ext in binary_extensions:
      return f"[Binary file '{os.path.basename(filepath)}' skipped]"

    if ext == ".pdf":
      try:
        res = subprocess.run(
            ["pdftotext", filepath, "-"],
            capture_output=True,
            text=True,
            timeout=10,
        )
        if res.returncode == 0:
          text = res.stdout
          if len(text) > max_chars:
            return text[:max_chars] + "\n...[TRUNCATED: PDF EXCEEDS 150,000 CHARACTERS]..."
          return text
      except Exception:
        pass
      return "[PDF File - Text extraction failed or pdftotext not available]"

    elif ext == ".docx":
      try:
        with zipfile.ZipFile(filepath) as docx:
          xml_content = docx.read("word/document.xml")
          root = ET.fromstring(xml_content)
          namespaces = {
              "w": (
                  "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
              )
          }
          text_elems = root.findall(".//w:t", namespaces)
          text = "".join([elem.text for elem in text_elems if elem.text])
          if len(text) > max_chars:
            return text[:max_chars] + "\n...[TRUNCATED: WORD DOC EXCEEDS 150,000 CHARACTERS]..."
          return text
      except Exception as e:
        return f"[Error reading Word file: {e}]"

    else:
      try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
          content = f.read(max_chars)
          if f.read(1):
            content += "\n...[TRUNCATED: FILE EXCEEDS 150,000 CHARACTERS]..."
          return content
      except Exception as e:
        return f"[Error reading file as text: {e}]"


  def evaluate(self) -> float:
    """Invokes the LLM Judge to evaluate the deliverables."""
    if GDPvalEnv._judge_client is None:
      try:
        from ..llm_client import VertexAIClient
        print("🧑‍⚖️ [GDPvalEnv] Initializing unified Gold Judge: gemini-3.1-pro-preview (location=global)...", flush=True)
        GDPvalEnv._judge_client = VertexAIClient(model_name="gemini-3.1-pro-preview", location="global", timeout=180)
      except Exception as e:
        print(f"⚠️ [GDPvalEnv] Failed to initialize Gold Judge gemini-3.1-pro-preview: {e}. Fallback to solver client.", flush=True)
        GDPvalEnv._judge_client = self.llm_client

    judge = GDPvalEnv._judge_client or self.llm_client
    if not judge:
      print("Warning: No LLM client available for evaluation. Returning 0.0")
      return 0.0


    # 1. Gather all files in the sandbox
    artifacts = {}
    for root, dirs, files in os.walk(self.sandbox_dir):
      # Prune large/hidden directories in-place to prevent walking into them
      dirs[:] = [
          d
          for d in dirs
          if d not in ("venv", ".git", "__pycache__", "node_modules")
          and not d.startswith(".")
      ]
      for file in files:
        full_path = os.path.join(root, file)
        rel_path = os.path.relpath(full_path, self.sandbox_dir)
        # Skip hidden files
        if file.startswith("."):
          continue
        artifacts[rel_path] = self._extract_text_from_file(full_path)

    if not artifacts:
      print("Warning: No artifacts found in sandbox. Score: 0.0")
      return 0.0

    # 2. Format artifacts for the prompt
    artifacts_content = ""
    for filename, content in artifacts.items():
      artifacts_content += f"<file_attachment>\n"
      artifacts_content += f"<file_name>{filename}</file_name>\n"
      artifacts_content += f"<file_content>\n{content}\n</file_content>\n"
      artifacts_content += f"</file_attachment>\n\n"

    # 3. Format rubrics
    rubrics_content = json.dumps(self.rubric_json, indent=2)

    # 4. Construct prompt
    prompt = JUDGE_PROMPT_TEMPLATE.format(
        artifacts_content=artifacts_content, rubrics_content=rubrics_content
    )

    # 5. Call LLM (We use temperature 0.0 for evaluation)
    # We use a system instruction to guide the judge
    system_instruction = (
        "You are a strict and precise grading assistant. Evaluate the"
        " deliverables exactly against the rubrics."
    )
    response = judge.generate(
        prompt, temperature=0.0, system_instruction=system_instruction
    )
    print(f"--- Raw Judge Response ---\n{response}\n--------------------------")

    # 6. Parse response
    try:
      # Extract JSON from <FINAL_RATINGS> tags, or fallback
      match = re.search(
          r"<FINAL_RATINGS>(.*?)</FINAL_RATINGS>", response, re.DOTALL
      )
      content = match.group(1).strip() if match else response

      # If it's a markdown block, strip it
      json_match = re.search(r"```json\s*(.*?)\s*```", content, re.DOTALL)
      if json_match:
        content = json_match.group(1).strip()

      ratings = None
      # Try standard json loads first with bracket cleaning
      try:
        cleaned_content = content
        bracket_match = re.search(r"(\[.*\])", content, re.DOTALL)
        if bracket_match:
          cleaned_content = bracket_match.group(1).strip()
        ratings = json.loads(cleaned_content)
      except Exception as je:
        print(f"⚠️ Standard parse failed: {je}. Trying robust extraction...")

      if ratings is None:
        # Robust extraction of individual criteria
        items = []
        starts = [m.start() for m in re.finditer(r'\{\s*"criterion \d+"\s*:', content)]
        for i, start_idx in enumerate(starts):
          end_idx = starts[i+1] if i + 1 < len(starts) else len(content)
          block = content[start_idx:end_idx].strip()
          if block.endswith(","):
            block = block[:-1].strip()
          if block.endswith("]"):
            block = block[:-1].strip()

          # Count open braces to fix missing closing braces
          open_braces = 0
          in_string = False
          escape = False
          for char in block:
            if escape:
              escape = False
              continue
            if char == '\\':
              escape = True
              continue
            if char == '"':
              in_string = not in_string
              continue
            if not in_string:
              if char == '{':
                open_braces += 1
              elif char == '}':
                open_braces -= 1

          if open_braces > 0:
            block += '}' * open_braces
          elif open_braces < 0:
            block = block[:open_braces]

          try:
            item = json.loads(block)
            items.append(item)
          except Exception as ie:
            print(f"Failed to parse block: {block} | Error: {ie}")
        ratings = items

      # Calculate score
      total_criteria = 0
      satisfied_criteria = 0
      
      if isinstance(ratings, dict):
        ratings = [ratings]
      elif not isinstance(ratings, list):
        ratings = []
        
      for item in ratings:
        if not isinstance(item, dict):
          continue
        for key, val in item.items():
          if "criterion" in key.lower():
            total_criteria += 1
            if isinstance(val, dict) and val.get("satisfied", False):
              satisfied_criteria += 1

      if total_criteria == 0:
        return 0.0
      score = satisfied_criteria / total_criteria
      print(
          f"Judge Evaluation: {satisfied_criteria}/{total_criteria} satisfied."
          f" Score: {score}"
      )
      return score

    except Exception as e:
      print(f"Error parsing judge response: {e}")
      print(f"Raw response:\n{response}")
      return 0.0

  # =========================================================================
  # Procedural Graph Hooks
  # =========================================================================

  def get_current_node_id(self, trajectory: List[str]) -> str:
    if not trajectory:
      return "START"

    # Find the most recent action.
    # To prevent getting stuck on old tool nodes when recent actions failed to parse,
    # we ONLY consider the last action in the trajectory.
    # If the last action failed to parse, we do NOT fallback to earlier actions; instead we return START.
    action_item = None
    for item in reversed(trajectory):
      if item.startswith("Action:"):
        action_item = item
        break  # Only look at the absolute most recent action!

    if action_item:
      action_content = action_item[len("Action:") :].strip()
      try:
        action_name, _, _ = env_base.parse_action_string(action_content)
        if action_name == "list_dir":
          return "Tool_list_dir"
        elif action_name == "search_files":
          return "Tool_search_files"
        elif action_name == "view_file":
          return "Tool_view_file"
        elif action_name == "read_excel":
          return "Tool_read_excel"
        elif action_name == "read_word":
          return "Tool_read_word"
        elif action_name == "read_pdf":
          return "Tool_read_pdf"
        elif action_name == "run_python":
          return "Tool_run_python"
        elif action_name == "write_file":
          return "Tool_write_file"
        elif action_name == "code_exec":
          return "Tool_code_exec"
        elif action_name == "submit":
          return "Tool_submit"
      except Exception:
        # If it failed to parse, return START to allow recovery guidelines
        return "START"

    return "START"
